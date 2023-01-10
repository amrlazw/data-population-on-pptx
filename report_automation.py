import psycopg2
from pptx import Presentation
from pptx.util import Inches

try:
    # Connect to the database
    conn = psycopg2.connect(
        host='localhost', 
        port=5301,
        database='mofdb',
        user='postgres', 
        password='qweasd'
    )
    print("-Successfully connected to the database!")
    # Create a cursor object
    cur = conn.cursor()

    # Fetch data from the database
    cur.execute("SELECT inisiatif, agensi, perbelanjaan FROM inisiatif_perbelanjaan")
    rows = cur.fetchall()

    # Create new presentation
    prs = Presentation()
    # Add a slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Create the table
    row_count = len(rows) + 1
    cols = 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(0.8)
    table = slide.shapes.add_table(row_count, cols, left, top, width, height).table
    
    #headers
    table.cell(0,0).text = 'Inisiatif'
    table.cell(0,1).text = 'Agensi'
    table.cell(0,2).text = 'Perbelanjaan (RM Juta)'
    
    #Insert data
    i = 1
    for row in rows:
        inisiatif, agensi, perbelanjaan = row
        table.cell(i, 0).text = inisiatif
        table.cell(i, 1).text = agensi
        table.cell(i, 2).text = str(perbelanjaan)
        i += 1

    # Save the PowerPoint presentation
    prs.save('Laporan_Bulanan_Generated.pptx')
    print("-Successfully created the pptx file!")

except psycopg2.Error as e:
    print(f'Error: {e}')
except Exception as e:
    print(f'Error: {e}')
finally:
    # Close the cursor and connection
    cur.close()
    conn.close()
