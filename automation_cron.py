import schedule
import time
import psycopg2
from pptx import Presentation
from pptx.util import Inches
def job():
    
# Connect to the database and fetch new data
    try:
        conn = psycopg2.connect(
            host='localhost', 
            port=5301,
            database='mofdb',
            user='postgres', 
            password='qweasd'
        )
        print("-Successfully connected to the database!")
        cur = conn.cursor()
        cur.execute("SELECT inisiatif, agensi, perbelanjaan FROM inisiatif_perbelanjaan")
        new_rows = cur.fetchall()
    finally:
        cur.close()
        conn.close()

    # Open the existing PowerPoint presentation
    try:
        prs = Presentation('Laporan_Bulanan_Generated.pptx')
        slide = prs.slides[0]
        table = slide.shapes[0].table

    except Exception as e:
        print(f'Error: {e}')

    # Determine the number of rows already in the table
    existing_row_count = len(table.rows)

    # Insert new data into the table
    i = existing_row_count
    for row in new_rows:
        inisiatif, agensi, perbelanjaan = row
        table.add_row()
        table.cell(i, 0).text = inisiatif
        table.cell(i, 1).text = agensi
        table.cell(i, 2).text = str(perbelanjaan)
        i += 1

    # Save the updated PowerPoint presentation
    prs.save('Laporan_Bulanan_Generated.pptx')
    print("-Successfully updated the pptx file!")
    
    pass

schedule.every(1).minutes.do(job)

while True:
    schedule.run_pending()
    time.sleep(1)
